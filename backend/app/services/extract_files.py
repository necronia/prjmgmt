"""업로드 파일 → 텍스트 추출. PDF/PPTX/DOCX/XLSX/텍스트는 로컬 파싱, 이미지는 Claude vision OCR."""
import io

from . import llm

IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp")


def _pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        txt = (page.extract_text() or "").strip()
        if txt:
            pages.append(f"[p.{i}]\n{txt}")
    return "\n\n".join(pages).strip()


def _pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            out.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
    return "\n\n".join(out).strip()


def _docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            out.append(f"[시트 {ws.title}]\n" + "\n".join(rows))
    return "\n\n".join(out).strip()


def extract_file(filename: str, content_type: str | None, data: bytes) -> tuple[str, str]:
    """(추출 텍스트, source_type) 반환. source_type ∈ {image, file}."""
    name = (filename or "").lower()
    ct = (content_type or "").lower()

    if ct.startswith("image/") or name.endswith(IMAGE_EXTS):
        return llm.ocr_image_bytes(data, ct or "image/png"), "image"
    if name.endswith(".pdf") or ct == "application/pdf":
        return _pdf(data), "file"
    if name.endswith(".pptx"):
        return _pptx(data), "file"
    if name.endswith(".docx"):
        return _docx(data), "file"
    if name.endswith(".xlsx"):
        return _xlsx(data), "file"
    # 텍스트류 (txt/md/csv/json/log 등)
    return data.decode("utf-8", errors="replace").strip(), "file"
